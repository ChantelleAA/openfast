"""
COMPREHENSIVE progress slides — past 2 weeks.
Every notebook, every plot, every intuition, every number.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu
from PIL import Image
import os

BASE = "/home/chantelle/Desktop/PhD Project/code/openfast"
DOUT = os.path.join(BASE, "demo_out")
OUT  = os.path.join(BASE, "meeting_slides.pptx")

DARK   = RGBColor(0x0D, 0x2B, 0x55)
MID    = RGBColor(0x1A, 0x6E, 0xAE)
LITE   = RGBColor(0xD0, 0xE8, 0xF5)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xF0, 0x8C, 0x00)
GREEN  = RGBColor(0x2E, 0xCC, 0x71)
GREY   = RGBColor(0x88, 0xAA, 0xCC)
RED    = RGBColor(0xE7, 0x4C, 0x3C)
YELLOW = RGBColor(0xF1, 0xC4, 0x0F)

W, H = Inches(13.33), Inches(7.5)
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── primitives ─────────────────────────────────────────────────────────────
def bg(s, c=None):
    f = s.background.fill; f.solid(); f.fore_color.rgb = c or DARK

def T(s, text, l, t, w, h, sz=16, bold=False, c=WHITE, al=PP_ALIGN.LEFT):
    b = s.shapes.add_textbox(l, t, w, h); b.word_wrap = True
    tf = b.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = al
    r  = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = c

def BL(s, items, l, t, w, h, title=None, sz=14, tsz=16):
    b = s.shapes.add_textbox(l, t, w, h); b.word_wrap = True
    tf = b.text_frame; tf.word_wrap = True
    if title:
        p = tf.paragraphs[0]; r = p.add_run()
        r.text = title; r.font.size = Pt(tsz)
        r.font.bold = True; r.font.color.rgb = ORANGE
    for item in items:
        p = tf.add_paragraph(); p.space_before = Pt(2)
        r = p.add_run()
        sub = item.startswith("   ") or item.startswith("  ")
        r.text = item if sub else (f"  {item}" if item=="" else f"  • {item}")
        r.font.size  = Pt(sz - (1 if sub else 0))
        r.font.color.rgb = GREY if item=="" else (LITE if sub else WHITE)

def HL(s, t, c=None):
    ln = s.shapes.add_connector(1, Inches(0.4), t, Inches(12.9), t)
    ln.line.color.rgb = c or MID; ln.line.width = Pt(1.2)

def IMG(s, path, l, t, w, h=None):
    if not os.path.exists(path): return
    if h: s.shapes.add_picture(path, l, t, w, h)
    else: s.shapes.add_picture(path, l, t, w)

def GIF(s, gif_path, l, t, w, h=None, frame=20):
    if not os.path.exists(gif_path): return
    im = Image.open(gif_path)
    try: im.seek(frame)
    except: pass
    tmp = "/tmp/_gf.png"; im.save(tmp)
    IMG(s, tmp, l, t, w, h)

def BOX(s, l, t, w, h, fill=None, border=None):
    r = s.shapes.add_shape(1, l, t, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = fill or MID
    r.line.color.rgb = border or WHITE; return r

def HDR(s, title, sub=None):
    T(s, title, Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.65),
      sz=24, bold=True)
    if sub:
        T(s, sub, Inches(0.4), Inches(0.78), Inches(12.5), Inches(0.38),
          sz=14, c=LITE)
    HL(s, Inches(0.78) if not sub else Inches(1.1))

# ══════════════════════════════════════════════════════════════════════════
# S1  Title
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
T(s,"Short-Term Extreme Value Prediction",
  Inches(0.8),Inches(1.3),Inches(11.7),Inches(1.1),sz=38,bold=True,al=PP_ALIGN.CENTER)
T(s,"for the IEA 15 MW Floating Offshore Wind Turbine",
  Inches(0.8),Inches(2.35),Inches(11.7),Inches(0.75),sz=26,c=LITE,al=PP_ALIGN.CENTER)
HL(s,Inches(3.2))
T(s,"ERA5 Environmental Data  ·  Joint Distribution Fitting  ·  IFORM Contours  ·  OpenFAST Simulation  ·  Gumbel & ACER Extreme Value Analysis",
  Inches(0.8),Inches(3.35),Inches(11.7),Inches(0.55),sz=13,c=RGBColor(0xAA,0xCC,0xEE),al=PP_ALIGN.CENTER)
T(s,"Chantelle  —  Supervisor Progress Update  |  April 2026",
  Inches(0.8),Inches(4.9),Inches(11.7),Inches(0.55),sz=17,c=ORANGE,al=PP_ALIGN.CENTER)
T(s,"Replicating: Chai et al. (2024) Ocean Engineering 306, 118120",
  Inches(0.8),Inches(5.6),Inches(11.7),Inches(0.45),sz=12,c=GREY,al=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
# S2  Big picture — what and why
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"The Big Picture: What Is This Work Trying to Answer?")
BL(s,[
  "Core question:  Is the IEA 15 MW floating wind turbine structurally safe under the worst 100-year sea conditions?",
  "",
  "To answer that, three problems must be solved in sequence:",
  "   1.  ENVIRONMENTAL PROBLEM — What is the worst realistic combination of wind + waves?",
  "        →  Build a 100-year environmental contour in (Uw, Hs, Tp) space",
  "   2.  SIMULATION PROBLEM — What loads does that combination produce on the turbine?",
  "        →  Run OpenFAST stochastic simulations at the extreme conditions",
  "   3.  STATISTICAL PROBLEM — What is the single worst load over the full return period?",
  "        →  Extrapolate extreme load distribution using Gumbel and ACER",
  "",
  "My work over the past 2 weeks covers ALL THREE stages:",
  "   Stage 1 — Validated the joint distribution model (Johannessen) and fitted it to ERA5 data",
  "   Stage 2 — Built the IFORM contour, selected ECs, set up OpenFAST simulation pipeline",
  "   Stage 3 — Implemented Gumbel and ACER extreme value prediction code",
  "",
  "PhD context — the ultimate goal is a 4th stage:",
  "   Stage 4 — Repeat the whole pipeline across time (1950–2025) to track how extremes change with climate",
],Inches(0.5),Inches(1.15),Inches(12.5),Inches(5.7),sz=14)

# ══════════════════════════════════════════════════════════════════════════
# S3  Paper overview — what Chai et al. 2024 did
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"Paper Overview: Chai et al. (2024)",
    "Ocean Engineering 306, 118120 — what they did, step by step")
steps = [
  ("Step 1","Build 3D environmental contour","Used Johannessen (2001) Northern North Sea stats. IFORM sphere (β=4.73) in (Uw, Hs, Tp) space gives Fig. 2."),
  ("Step 2","Slice at 4 hub wind speeds","Fixed Vhub = 6.98 · 10.59 · 25.0 · 57.64 m/s → 2D (Hs,Tp) loops at each speed (Fig. 3)."),
  ("Step 3","Pick worst (Hs, Tp) per loop","Run OpenFAST at several points on each loop → find maximum structural response → that is the design EC."),
  ("Step 4","50 stochastic simulations per EC","JONSWAP waves, different random seeds, 4200 s each (600 s transient + 3600 s analysis)."),
  ("Step 5","Extreme value prediction","Gumbel block-maxima and ACER peak-exceedance applied to mooring tension, tower shear, bending moment."),
  ("Step 6","Safety check","Mooring: max tension < 60% MBS (API-RP-2SK). Tower: within Eurocode 3 S355 allowables. Both pass ✓"),
]
y = Inches(1.2)
for tag, title, body in steps:
    BOX(s, Inches(0.4), y, Inches(1.1), Inches(0.88), fill=MID)
    T(s, tag, Inches(0.4), y+Inches(0.2), Inches(1.1), Inches(0.5),
      sz=13, bold=True, al=PP_ALIGN.CENTER)
    T(s, title, Inches(1.6), y, Inches(11.2), Inches(0.38), sz=15, bold=True, c=ORANGE)
    T(s, body,  Inches(1.8), y+Inches(0.37), Inches(11.0), Inches(0.48), sz=13, c=LITE)
    y += Inches(1.0)

# ══════════════════════════════════════════════════════════════════════════
# S4  Surprising finding
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"Key Insight from the Paper: Why the Parked Case Is Worst",
    "EC5 (57.64 m/s, turbine parked) drives the largest tower loads — counterintuitively")
BL(s,[
  "You might expect the operating turbine (blades spinning, converting wind to electricity) to be under more stress.",
  "But the parked turbine (blades feathered, no rotation) actually produces the LARGEST tower loads.",
  "",
  "Why?  When the turbine is operating, the spinning rotor acts like a giant sail that 'catches' the wind",
  "   smoothly and converts it to torque — the aerodynamic loads are distributed over the rotor disk.",
  "   The controller also actively limits thrust by pitching blades.",
  "",
  "   When the turbine is PARKED, the blades are feathered (edge-on to the wind) but the wind still",
  "   hits the tower and platform directly with nothing absorbing it.",
  "   The drag on the tower + platform becomes the dominant structural load.",
  "",
  "At EC5: Vhub = 57.64 m/s, Hs = 14.33 m, Tp = 14.78 s",
  "   Tower base shear force:     5.15 × 10³ kN    (maximum across all ECs)",
  "   Tower base bending moment:  7.82 × 10⁵ kN·m  (maximum across all ECs)",
  "   Max mooring tension: 4.74 × 10³ kN  (safety factor 4.70 — well above minimum 1.67 ✓)",
  "",
  "This shapes our simulation strategy: EC5 (parked, 100-year wind) is the critical design load case.",
],Inches(0.5),Inches(1.15),Inches(12.5),Inches(5.6),sz=14)

# ══════════════════════════════════════════════════════════════════════════
# S5  ERA5 data
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"Data Foundation: ERA5 Reanalysis  (1950–2025)",
    "75 years of hourly metocean data — the input to everything else")
BL(s,[
  "Source: ECMWF ERA5 global atmospheric reanalysis, downloaded via the Copernicus CDS API",
  "Variables extracted at each grid point:",
  "   u10, v10    →  10 m wind components  →  W = √(u²+v²)  (m/s)",
  "   u100, v100  →  100 m wind components  (used for validation + sensitivity)",
  "   swh         →  significant wave height Hs  (m)",
  "   pp1d        →  peak spectral wave period Tp  (s)",
  "",
  "Data summary at validation point (56.5°N, 3.5°E — North of Scotland):",
  "   Total records:  666,216 hourly observations  (1950–2025, 76 years)",
  "   Mean wind W:    8.35 m/s  ·  Weibull shape β=2.35, scale α=9.42 m/s",
  "   100-year wind (10 m, 1-h):  28.7 m/s  (from fitted distribution)",
  "",
  "6 Irish Sea / North Sea buoy locations extracted (M1–M6):",
  "   M1: 53.13°N, 11.20°W   M2: 53.48°N, 5.43°W   M3: 51.22°N, 10.55°W",
  "   M4: 55.00°N, 10.00°W   M5: 51.69°N, 6.70°W   M6: 53.06°N, 15.93°W",
  "",
  "Current demo subset used for preliminary EVA: 1950–1981  (31.2 years, 375 files)",
  "   Enough for 32 annual maxima → meaningful Gumbel/GEV/POT comparison",
],Inches(0.5),Inches(1.15),Inches(7.5),Inches(5.6),sz=14)
IMG(s,os.path.join(DOUT,"fig_timeseries_vhub.png"),Inches(7.8),Inches(1.15),Inches(5.3))
T(s,"Hub wind speed time series\n(10 m extrapolation, 150 m)",
  Inches(7.8),Inches(5.5),Inches(5.3),Inches(0.5),sz=11,c=GREY,al=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
# S6  Hub-height conversion formula + intuition
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"Hub-Height Wind Extrapolation: Formula and Intuition",
    "ERA5 measures wind at 10 m. The turbine hub is at 150 m. Two corrections needed.")
BL(s,[
  "Physical intuition: Wind is slower near the sea surface due to friction.",
  "   It increases with height — this is the wind SHEAR PROFILE.",
  "   A 150 m hub sits in significantly faster wind than the surface.",
  "",
  "Step 1 — Height correction (power-law profile, IEC 61400-1):",
  "   U_hub = U_ref × (z_hub / z_ref)^α      with α = 0.11 (offshore standard)",
  "   From 10 m:   U_hub = U_10  × (150/10)^0.11  =  U_10  × 1.369",
  "   From 100 m:  U_hub = U_100 × (150/100)^0.11 =  U_100 × 1.046",
  "",
  "Step 2 — Averaging period (1-hour ERA5 → 10-minute turbine standard):",
  "   V_hub (10-min) = U_hub (1-h) × 1.10",
  "   A 10-min average is higher than 1-h because shorter windows capture stronger gusts.",
  "",
  "Combined from 10 m:  V_hub = U_10 × 1.10 × 1.369  ≈  U_10 × 1.506",
  "   Example: U_10 = 38.9 m/s  →  V_hub = 57.6 m/s  ✓  (matches Chai et al. EC5)",
  "",
  "Why larger α means more amplification:",
  "   α = 0.05 (very flat profile): (150/10)^0.05 = 1.145",
  "   α = 0.11 (offshore standard): (150/10)^0.11 = 1.369",
  "   α = 0.20 (onshore, rougher):  (150/10)^0.20 = 1.741",
  "   → 10 m extrapolation is HIGHLY SENSITIVE to α choice (large vertical jump)",
],Inches(0.5),Inches(1.15),Inches(7.2),Inches(5.6),sz=14)
# formula boxes
for i,(lbl,eq,val) in enumerate([
    ("10 m → 150 m","U₁₀ × 1.369","×1.10 → V_hub"),
    ("100 m → 150 m","U₁₀₀ × 1.046","×1.10 → V_hub"),
]):
    BOX(s,Inches(7.5),Inches(1.5+i*2.3),Inches(5.5),Inches(2.0),
        fill=RGBColor(0x10,0x32,0x60))
    T(s,lbl,Inches(7.6),Inches(1.6+i*2.3),Inches(5.3),Inches(0.45),
      sz=14,bold=True,c=ORANGE,al=PP_ALIGN.CENTER)
    T(s,eq,Inches(7.6),Inches(2.05+i*2.3),Inches(5.3),Inches(0.5),
      sz=18,bold=True,al=PP_ALIGN.CENTER)
    T(s,val,Inches(7.6),Inches(2.55+i*2.3),Inches(5.3),Inches(0.4),
      sz=13,c=GREEN,al=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
# S7  10 m vs 100 m comparison plots
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"10 m vs 100 m Extrapolation: Which Is More Reliable?",
    "Internal validation against actual ERA5 100 m wind reveals systematic bias in 10 m method")
BL(s,[
  "Two hub-height series computed from the same ERA5 data:",
  "   From 10 m:   mean V_hub = 12.37 m/s",
  "   From 100 m:  mean V_hub = 11.55 m/s",
  "   Mean absolute difference: 0.96 m/s  (8.66% of the 100 m value)",
  "",
  "Internal validation — extrapolate 10 m ONLY UP TO 100 m, compare to actual ERA5 100 m:",
  "   Extrapolated 10→100 m mean: 10.75 m/s",
  "   Actual ERA5 100 m mean:     10.04 m/s",
  "   Bias: +0.71 m/s  ·  MAE: 0.84 m/s  ·  RMSE: 0.97 m/s  ·  MAPE: 9.76%",
  "",
  "Conclusion: 10 m power-law already overpredicts at 100 m",
  "   → likely also overpredicts at 150 m",
  "   → 100 m based extrapolation is the more trustworthy estimate",
  "   → 10 m based extrapolation is kept as the paper-consistent BASELINE",
],Inches(0.4),Inches(1.15),Inches(5.5),Inches(4.5),sz=13)
IMG(s,os.path.join(DOUT,"fig_monthly_mean_vhub_compare.png"),
    Inches(5.9),Inches(1.1),Inches(3.6))
IMG(s,os.path.join(DOUT,"fig_monthly_mean_ws100_validation.png"),
    Inches(9.65),Inches(1.1),Inches(3.5))
IMG(s,os.path.join(DOUT,"fig_extrapolation_means_bar.png"),
    Inches(5.9),Inches(4.55),Inches(3.6))
IMG(s,os.path.join(DOUT,"fig_extrapolation_gap.png"),
    Inches(9.65),Inches(4.55),Inches(3.5))
for lbl, x in [("Monthly means comparison",5.9),("100 m validation",9.65),("Mean bar chart",5.9),("Extrapolation gap",9.65)]:
    pass  # labels already in figures

# ══════════════════════════════════════════════════════════════════════════
# S8  Alpha sensitivity
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"Sensitivity to Power-Law Exponent α",
    "10 m extrapolation is much more sensitive to α than 100 m — quantified here")
BL(s,[
  "α (shear exponent) controls how fast wind increases with height.",
  "   Offshore seas: α ≈ 0.10–0.12  (smooth surface, weak shear)",
  "   Coastal/nearshore: α ≈ 0.14–0.16",
  "   Onshore/rough: α ≈ 0.20–0.25",
  "",
  "Key result from sensitivity analysis:",
  "   At α = 0.11: 10 m gives 12.37 m/s · 100 m gives 11.55 m/s",
  "   As α increases, 10 m based estimate rises MUCH faster",
  "   Because: 10 m spans 10→150 m (15× height ratio) vs 100 m spans 100→150 m (1.5× ratio)",
  "",
  "This is why the 10 m extrapolation is high-risk:",
  "   A small error in α (which is often region- and season-dependent)",
  "   gets amplified over the full 140 m vertical span",
  "",
  "For thesis: will use α = 0.11 as paper-consistent baseline,",
  "   carry α sensitivity as formal uncertainty quantification",
],Inches(0.4),Inches(1.15),Inches(5.8),Inches(5.5),sz=14)
IMG(s,os.path.join(DOUT,"fig_alpha_sensitivity.png"),
    Inches(6.2),Inches(1.1),Inches(6.9))

# ══════════════════════════════════════════════════════════════════════════
# S9  Annual maxima & time series comparison
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"Annual Maxima and Time Series — Input to Extreme Value Analysis")
IMG(s,os.path.join(DOUT,"fig_timeseries_vhub_compare.png"),
    Inches(0.4),Inches(1.15),Inches(6.1))
IMG(s,os.path.join(DOUT,"fig_annual_maxima.png"),
    Inches(6.7),Inches(1.15),Inches(6.4))
IMG(s,os.path.join(DOUT,"fig_annual_maxima_compare.png"),
    Inches(0.4),Inches(4.35),Inches(6.1))
BL(s,[
  "Top left: full hub-wind time series  (10 m vs 100 m paths track closely in shape,",
  "   differ by ~0.96 m/s mean offset)",
  "",
  "Top right: annual maxima used as input to Gumbel/GEV fitting",
  "   32 years of data → 32 block maxima",
  "   Scatter reflects natural interannual variability",
  "",
  "Bottom left: annual maxima comparison between 10 m and 100 m paths",
  "   Similar year-to-year pattern — same storm events drive both",
  "   10 m consistently higher due to larger extrapolation factor",
],Inches(6.7),Inches(4.35),Inches(6.4),Inches(2.9),sz=13)

# ══════════════════════════════════════════════════════════════════════════
# S10  Five EVA methods
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"Five Extreme Value Methods — Quantifying Model-Form Uncertainty",
    "Applied to hub-height wind speed from ERA5 subset (1950–1981, 32 annual maxima)")
BL(s,[
  "Why compare five methods?  Any single EVA method embeds assumptions.",
  "   Comparing them quantifies model-form uncertainty — how much the estimate depends on model choice.",
],Inches(0.4),Inches(1.15),Inches(12.5),Inches(0.7),sz=14)
methods = [
  ("Gumbel\nblock-max","F(x)=exp(-exp(-(x-μ)/β))","Simple, transparent, used in paper. CONSERVATIVE.","36.31"),
  ("GEV\nblock-max","F(x)=exp(-(1+ξ(x-μ)/σ)^(-1/ξ))","Generalises Gumbel (ξ=0). More flexible.","32.90"),
  ("Weibull\nblock-max","F(x)=1-exp(-(x/α)^β)","Practical engineering benchmark for positive vars.","32.88"),
  ("Lognormal\nblock-max","ln(X) ~ Normal(μ,σ)","Benchmark for positive skewed variables.","33.68"),
  ("POT/GPD","P(Y>y|X>u)=(1+ξy/σᵤ)^(-1/ξ)","Uses more tail data. Threshold u=21.89 m/s, n=746 peaks.","33.16"),
]
for i,(nm,formula,note,rl100) in enumerate(methods):
    y = Inches(1.95+i*0.95)
    rb = RGBColor(0x12,0x35,0x6A) if i%2==0 else DARK
    BOX(s,Inches(0.4),y,Inches(1.7),Inches(0.85),fill=rb,border=MID)
    T(s,nm,Inches(0.45),y+Inches(0.1),Inches(1.6),Inches(0.7),sz=11,bold=True,
      c=ORANGE if i==0 else WHITE,al=PP_ALIGN.CENTER)
    BOX(s,Inches(2.2),y,Inches(3.9),Inches(0.85),fill=rb,border=MID)
    T(s,formula,Inches(2.25),y+Inches(0.15),Inches(3.8),Inches(0.6),sz=12,c=LITE)
    BOX(s,Inches(6.2),y,Inches(5.5),Inches(0.85),fill=rb,border=MID)
    T(s,note,Inches(6.25),y+Inches(0.15),Inches(5.4),Inches(0.6),sz=12,c=WHITE)
    BOX(s,Inches(11.8),y,Inches(1.3),Inches(0.85),fill=rb,border=MID)
    T(s,rl100+" m/s",Inches(11.82),y+Inches(0.2),Inches(1.26),Inches(0.5),
      sz=13,bold=True,c=ORANGE if i==0 else GREEN,al=PP_ALIGN.CENTER)

T(s,"100-yr RL",Inches(11.8),Inches(1.95-0.4),Inches(1.3),Inches(0.38),
  sz=11,bold=True,c=GREY,al=PP_ALIGN.CENTER)
T(s,"Key finding: Gumbel is 9–10% MORE CONSERVATIVE than the other methods (36.31 vs ~33 m/s).",
  Inches(0.4),Inches(6.75),Inches(9),Inches(0.45),sz=13,c=ORANGE)
T(s,"This is exactly why the paper pairs Gumbel with ACER — one conservative, one efficient.",
  Inches(0.4),Inches(7.1),Inches(9),Inches(0.38),sz=13,c=LITE)

# ══════════════════════════════════════════════════════════════════════════
# S11  Return level + POT plots
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"Return Level Comparison and POT Threshold Diagnostic")
IMG(s,os.path.join(DOUT,"fig_return_levels.png"),
    Inches(0.4),Inches(1.15),Inches(6.3))
IMG(s,os.path.join(DOUT,"fig_pot_mrl.png"),
    Inches(6.9),Inches(1.15),Inches(6.2))
BL(s,[
  "Left — Return level plot:",
  "   All 5 methods plotted on the same axes",
  "   Gumbel diverges upward at long return periods",
  "   GEV, Weibull, Lognormal, POT cluster tightly — they agree",
  "   Spread at 100 years quantifies model-form uncertainty",
  "",
  "Right — Mean Residual Life (MRL) plot for POT threshold selection:",
  "   x-axis: candidate threshold u",
  "   y-axis: mean excess (average amount by which exceedances exceed u)",
  "   Theory: if data follows GPD above u, MRL should be LINEAR in u",
  "   The near-linear region identifies a valid threshold range",
  "   Chosen threshold: u = 21.89 m/s  (q=0.95), giving 746 declustered peaks",
  "   GPD shape parameter ξ = −0.327  (upper tail is bounded)",
],Inches(0.4),Inches(4.5),Inches(12.5),Inches(2.75),sz=13)

# ══════════════════════════════════════════════════════════════════════════
# S12  Joint distribution — the 3 variables
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"Joint Distribution Model: (Uw, Hs, Tp) Are Not Independent",
    "Strong physical coupling — wind drives waves, waves set the period")
BL(s,[
  "The three environmental variables are physically linked:",
  "   High winds → large waves → long-period swell",
  "   Must model them JOINTLY — not independently",
  "",
  "Hierarchical conditional factorisation (Rosenblatt transform):",
  "   f(Uw, Hs, Tp)  =  f(Uw)  ×  f(Hs | Uw)  ×  f(Tp | Hs, Uw)",
  "",
  "Three distributions fitted:",
],Inches(0.4),Inches(1.15),Inches(5.8),Inches(3.8),sz=14)
rows = [
  ("Variable","Distribution","Parameters depend on","Source"),
  ("Uw (wind)","Marginal Weibull","—  (fitted to all data)","ERA5 10 m wind"),
  ("Hs | Uw","Conditional Weibull","Scale & shape linear in Uw","ERA5 swh"),
  ("Tp | Hs, Uw","Conditional Lognormal","μ & σ as power-law in Hs","ERA5 pp1d"),
]
widths=[Inches(1.5),Inches(2.2),Inches(2.8),Inches(1.8)]
xs=[Inches(0.4),Inches(2.0),Inches(4.3),Inches(7.2)]
for ri,row in enumerate(rows):
    yrow=Inches(4.2+ri*0.68)
    rb=MID if ri==0 else (RGBColor(0x12,0x35,0x6A) if ri%2==1 else DARK)
    for val,x,w in zip(row,xs,widths):
        BOX(s,x,yrow,w,Inches(0.63),fill=rb,border=WHITE)
        T(s,val,x+Inches(0.04),yrow+Inches(0.08),w-Inches(0.08),Inches(0.55),
          sz=12,bold=(ri==0),al=PP_ALIGN.CENTER,
          c=ORANGE if ri==0 else WHITE)
BL(s,[
  "Johannessen (2001) parameters  (benchmark):",
  "   Uw Weibull:  α=8.43, β=1.71",
  "   α_Hs(W) = 1.8 + 0.100·W^1.322",
  "   β_Hs(W) = 2.0 + 0.135·W",
  "   Tp model: μ_Tp = (4.883+2.68·Hs^0.529)·(1−0.19·Δ)",
  "",
  "Our ERA5 fit  (56.5°N):",
  "   Uw Weibull:  α=9.42, β=2.35",
  "   α_Hs(W) = 0.488·W^0.870 + 0.829",
  "   β_Hs(W) = 0.011·W^2.106 + 0.995",
],Inches(6.3),Inches(1.2),Inches(6.8),Inches(5.5),sz=13)

# ══════════════════════════════════════════════════════════════════════════
# S13  Weibull wind fit plot
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"Fitting the Marginal Wind Distribution: Weibull MLE",
    "Method-of-moments initialisation → scipy MLE → validated against empirical CDF")
BL(s,[
  "Procedure:",
  "   1. Compute wind speed: W = √(u10²+v10²) from ERA5 components",
  "   2. Method-of-moments: estimate shape β and scale α from sample mean and CV",
  "   3. Maximum likelihood estimation (scipy weibull_min.fit) refines parameters",
  "   4. Validate: overlay Weibull PDF on histogram, CDF on empirical CDF",
  "",
  "ERA5 results at 56.5°N, 3.5°E:",
  "   Shape β = 2.35  (heavier than Johannessen's 1.71 — more Rayleigh-like)",
  "   Scale α = 9.42 m/s  (slightly higher than Johannessen's 8.43 m/s)",
  "",
  "Fitted 100-year return level: 28.7 m/s  (10 m, 1-h)",
  "Johannessen (Northern North Sea): ~39 m/s  (harsher, deeper water)",
  "",
  "The difference is real — our sites are in the milder Irish Sea,",
  "   not the open, deep Northern North Sea of the Johannessen dataset",
],Inches(0.4),Inches(1.15),Inches(5.8),Inches(5.5),sz=14)
IMG(s,os.path.join(BASE,"weibull_wind_fit.png"),Inches(6.3),Inches(1.1),Inches(6.8))

# ══════════════════════════════════════════════════════════════════════════
# S14  Conditional Hs parameters
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"Conditional Wave Height: Hs | Uw",
    "Weibull scale and shape as smooth functions of wind speed — fitted from 14 bins")
BL(s,[
  "For each wind speed bin, Hs follows a conditional Weibull:",
  "   F(h|w) = 1 − exp[−(h / α_Hs(w))^β_Hs(w)]",
  "",
  "Fitting procedure:",
  "   Bin ERA5 data into 14 wind speed intervals  (0–2, 2–4, ..., 24–26 m/s)",
  "   Minimum 19,000 data points per bin  (dense enough for stable fit)",
  "   Fit Weibull to Hs in each bin → get one (α, β) pair per bin",
  "   Fit SMOOTH power-law functions through the bin estimates:",
  "      α_Hs(W) = 0.488·W^0.870 + 0.829",
  "      β_Hs(W) = 0.011·W^2.106 + 0.995",
  "",
  "Physical interpretation:",
  "   As wind speed increases, both scale (typical wave height) and shape",
  "   (spread of the distribution) increase — bigger and more variable waves",
  "",
  "The smooth functions are what get plugged into the IFORM contour generator",
],Inches(0.4),Inches(1.15),Inches(5.8),Inches(5.5),sz=14)
IMG(s,os.path.join(BASE,"Hs_smooth_params.png"),Inches(6.3),Inches(1.1),Inches(6.8))

# ══════════════════════════════════════════════════════════════════════════
# S15  Conditional Tp lognormal
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"Conditional Peak Period: Tp | Uw, Hs",
    "Lognormal distribution — parameters fitted in (Uw, Hs) double-bins")
BL(s,[
  "Peak period Tp ~ Lognormal  given (Uw, Hs):",
  "   ln(Tp) ~ Normal(μ_ln, σ_ln)",
  "",
  "Double-bin fitting:",
  "   First bin by Uw (6 bins), then by Hs within each Uw bin (6 sub-bins)",
  "   Minimum 30 points per (Uw, Hs) cell",
  "   In each cell: fit μ_ln and σ_ln of log(Tp)",
  "   Then smooth these as power-law functions of Hs:",
  "      μ_ln(Hs) = a_μ + b_μ · Hs^c_μ",
  "      σ_ln(Hs) = a_σ + b_σ · Hs^c_σ",
  "",
  "Key values:",
  "   W=10 m/s, Hs=5 m  →  median Tp = 13.1 s",
  "   W=15 m/s, Hs=8 m  →  median Tp = 15.2 s",
  "   W=25 m/s, Hs=13 m →  median Tp = 16.2 s",
  "",
  "Note: extreme (W, Hs) combinations at high wind + low wave height",
  "   can cause parameter overflow — flagged and filtered",
],Inches(0.4),Inches(1.15),Inches(5.8),Inches(5.5),sz=14)
IMG(s,os.path.join(BASE,"lognormal_Tp_params.png"),Inches(6.3),Inches(1.1),Inches(6.8))

# ══════════════════════════════════════════════════════════════════════════
# S16  IFORM theory
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"IFORM: How the Return Period Becomes a 3D Surface",
    "Inverse First-Order Reliability Method — the mathematical bridge from statistics to design conditions")
BL(s,[
  "Step 1 — Map to standard normal space (Rosenblatt transform):",
  "   U₁ = Φ⁻¹(F_W(Uw))",
  "   U₂ = Φ⁻¹(F_{Hs|W}(Hs | Uw))",
  "   U₃ = Φ⁻¹(F_{Tp|Hs,W}(Tp | Hs, Uw))",
  "   In (U₁,U₂,U₃) space, all distributions become standard normal",
  "",
  "Step 2 — The 100-year event maps to a sphere of radius β:",
  "   β = Φ⁻¹(1 − 1/(T · N_obs/yr))",
  "   T=100 yr, N_obs=8766/yr  →  β = 4.727",
  "   Every point on the sphere has the SAME return period — just different combinations",
  "",
  "Step 3 — Parametrise the sphere with (θ, φ) and transform back:",
  "   u₁ = β·cos(θ)                     →  Φ⁻¹ of F_W  →  Uw",
  "   u₂ = β·sin(θ)·cos(φ)  + Uw  →  Φ⁻¹ of F_{Hs|W}  →  Hs",
  "   u₃ = β·sin(θ)·sin(φ)  + Uw, Hs →  Φ⁻¹ of F_{Tp|Hs,W}  →  Tp",
  "",
  "Result: a 3D surface in (Uw, Hs, Tp) — every point is a 100-year event",
  "Discretisation: θ ∈ [0,π] × 50 points,  φ ∈ [0,2π] × 100 points  =  5000 surface points",
],Inches(0.4),Inches(1.15),Inches(7.5),Inches(5.6),sz=14)
# sphere sketch
BOX(s,Inches(7.9),Inches(1.4),Inches(5.1),Inches(5.7),
    fill=RGBColor(0x08,0x1E,0x40),border=MID)
T(s,"Standard Normal\nSpace (U₁,U₂,U₃)",Inches(8.0),Inches(1.5),Inches(4.9),Inches(0.6),
  sz=13,bold=True,c=ORANGE,al=PP_ALIGN.CENTER)
T(s,"Sphere of radius β=4.727\n\nEvery point on\nthis sphere\n= 100-year event",
  Inches(8.1),Inches(2.2),Inches(4.7),Inches(2.5),sz=14,c=WHITE,al=PP_ALIGN.CENTER)
T(s,"Rosenblatt\ntransform\n↕\n(Uw, Hs, Tp)\nphysical space",
  Inches(8.1),Inches(4.8),Inches(4.7),Inches(2.0),sz=13,c=GREEN,al=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
# S17  Johannessen 3D contour
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"Johannessen 50-Year Contour — Benchmark Validated",
    "Reproduced from parametric model. Max values match published results within 3%.")
BL(s,[
  "Johannessen (2001) Northern North Sea, 1973–1999:",
  "   K1Y = 8760 obs/yr  ·  β = 4.584 for 50-year return",
  "",
  "Contour surface results:",
  "   Max Uw = 37.8 m/s (10 m, 1-h)",
  "   Hub conversion: 37.8 × 1.10 × (150/10)^0.11 = 56.0 m/s",
  "   Paper reports: 57.64 m/s  →  3% difference ✓",
  "   Max Hs = 15.2 m  ·  Max Tp = 23.8 s",
  "",
  "Surface shape (Fig. 2 of paper):",
  "   Low winds (left corner) → small waves, short periods",
  "   High winds (tall peak) → large waves, long swell",
  "   The surface 'wraps' around the IFORM reliability sphere",
  "",
  "Axes are intentionally inverted (Hs axis flips, Tp axis flips)",
  "   to match the orientation of Fig. 2 in Chai et al. 2024",
],Inches(0.4),Inches(1.15),Inches(5.7),Inches(5.5),sz=14)
IMG(s,os.path.join(BASE,"johannessen_3D_matched.png"),Inches(6.2),Inches(0.9),Inches(6.9))

# ══════════════════════════════════════════════════════════════════════════
# S18  Rotating GIF — both side by side
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"360° Rotating Contour Animations",
    "120 frames · 20 fps · matplotlib FuncAnimation + PillowWriter — full azimuth sweep")
GIF(s,os.path.join(DOUT,"rotating_3d_environmental_contour_johannesen.gif"),
    Inches(0.4),Inches(1.15),Inches(5.95),frame=25)
GIF(s,os.path.join(DOUT,"rotating_3d_environmental_contour.gif"),
    Inches(7.0),Inches(1.15),Inches(5.95),frame=25)
T(s,"Johannessen model — 50-year contour\n(benchmark, Northern North Sea)",
  Inches(0.4),Inches(5.75),Inches(5.95),Inches(0.6),sz=12,c=GREY,al=PP_ALIGN.CENTER)
T(s,"ERA5-fitted model — 100-year contour\n(our data, Irish Sea buoy M1)",
  Inches(7.0),Inches(5.75),Inches(5.95),Inches(0.6),sz=12,c=GREY,al=PP_ALIGN.CENTER)
HL(s,Inches(6.65),c=RGBColor(0x44,0x66,0x88))
T(s,"What the animation reveals that a static image cannot:",
  Inches(0.4),Inches(6.5),Inches(5.5),Inches(0.38),sz=13,bold=True,c=ORANGE)
T(s,"the surface is a thin shell at low winds, opens into a broad dome at high winds — confirming the physical meaning of the joint distribution",
  Inches(0.4),Inches(6.85),Inches(12.5),Inches(0.45),sz=12,c=LITE)

# ══════════════════════════════════════════════════════════════════════════
# S19  ERA5 contour + comparison
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"ERA5-Fitted 100-Year Contour — Our Own Data",
    "Fitted independently using 75 years of ERA5 at 6 buoy locations")
BL(s,[
  "All distributions fitted from ERA5 data — not Johannessen's parameters",
  "",
  "M1 (53.1°N, 11.2°W) 100-year results:",
  "   β = 4.727  (100 yr × 8766 obs/yr)",
  "   Max Uw = 29.5 m/s  ·  Max Hs = 8.49 m  ·  Max Tp = 29.1 s",
  "",
  "Why lower than Johannessen?",
  "   Johannessen used the open Northern North Sea (water depth >200 m)",
  "   Our M1–M6 locations are in the shallower Irish Sea / Celtic Sea",
  "   These are genuinely less severe metocean conditions",
  "",
  "Hs comparison at max Uw:",
  "   Johannessen 50-yr:  Hs_max = 15.2 m",
  "   Our ERA5 100-yr:    Hs_max = 8.49 m",
  "   Physically consistent — milder sea, even at a LONGER return period",
  "",
  "The 6 buoy locations give slightly different contour shapes,",
  "   reflecting regional variation in the Irish/North Sea climatology",
],Inches(0.4),Inches(1.15),Inches(5.8),Inches(5.5),sz=14)
IMG(s,os.path.join(BASE,"contour_surface_3D.png"),Inches(6.3),Inches(0.9),Inches(6.8))

# ══════════════════════════════════════════════════════════════════════════
# S20  2D contour slices
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"2D Hs–Tp Loops: Slicing the 3D Contour at Fixed Wind Speed",
    "Each coloured loop = one hub wind speed = one 'equally dangerous' family of wave conditions")
BL(s,[
  "Procedure:",
  "   Fix V_hub  →  convert to Uw (10 m, 1-h) using Eq. 9",
  "   U_w = V_hub / (1.10 × (150/10)^0.11)",
  "   Find the angle θ on the IFORM sphere corresponding to that Uw",
  "   Sweep φ from 0 to 2π  →  trace a closed curve in (Hs, Tp) space",
  "",
  "Hub speeds used:  6.98 · 10.59 · 25.0 · 57.64 m/s",
  "   (cut-in, rated, cut-out, 100-year survival)",
  "",
  "Worst-case = point with MAX Hs on each loop",
  "   Physically: biggest waves produce the largest structural loads",
  "   This is the (Hs, Tp) pair fed into OpenFAST for that wind speed",
  "",
  "Note: 57.64 m/s is right at the edge of the contour sphere",
  "   (Uw = 38.9 m/s is near the maximum of the Weibull distribution)",
],Inches(0.4),Inches(1.15),Inches(5.5),Inches(4.4),sz=14)
IMG(s,os.path.join(BASE,"johannessen_2D_fixed.png"),Inches(0.4),Inches(4.55),Inches(6.0))
IMG(s,os.path.join(BASE,"M1_2D_fixed_hub_speeds.png"),Inches(6.7),Inches(0.85),Inches(6.4))
T(s,"Johannessen benchmark (bottom left)                         ERA5 M1 fitted (right)",
  Inches(0.4),Inches(7.05),Inches(12.5),Inches(0.38),sz=11,c=GREY)

# ══════════════════════════════════════════════════════════════════════════
# S21  EC table
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"Selected Environmental Conditions (ECs)",
    "One worst-case (Hs, Tp) per wind speed — the simulation inputs")
T(s,"For each hub wind speed, the point with maximum Hs on the 100-year contour loop is selected.",
  Inches(0.4),Inches(1.15),Inches(12.5),Inches(0.42),sz=14,c=LITE)
cols   =["EC","V_hub (m/s)\n10-min, 150 m","Uw (m/s)\n1-h, 10 m","Hs (m)","Tp (s)","Turbine state"]
widths =[Inches(0.75),Inches(1.9),Inches(1.7),Inches(1.3),Inches(1.3),Inches(2.3)]
xs     =[Inches(0.4),Inches(1.25),Inches(3.25),Inches(5.05),Inches(6.45),Inches(7.85)]
for col,x,w in zip(cols,xs,widths):
    BOX(s,x,Inches(1.65),w,Inches(0.58),fill=MID)
    T(s,col,x+Inches(0.04),Inches(1.67),w-Inches(0.08),Inches(0.58),
      sz=11,bold=True,al=PP_ALIGN.CENTER)
rows_ec=[
  ("EC1","6.98","4.71","~3–4","~10–11","Operation (below rated)"),
  ("EC2","10.59","7.15","~5–6","~11–12","Operation (rated)"),
  ("EC3","25.0","16.87","~9–11","~13–14","Operation (above rated)"),
  ("EC4","57.64","38.9","~13–15","~14–16","Parked — SURVIVAL"),
]
for ri,row in enumerate(rows_ec):
    yr=Inches(2.33+ri*0.7)
    rb=RGBColor(0x12,0x35,0x6A) if ri%2==0 else DARK
    for val,x,w in zip(row,xs,widths):
        BOX(s,x,yr,w,Inches(0.65),fill=rb,border=MID)
        T(s,val,x+Inches(0.04),yr+Inches(0.07),w-Inches(0.08),Inches(0.57),
          sz=12,c=ORANGE if val==row[0] else (RED if "SURVIVAL" in val else WHITE),
          al=PP_ALIGN.CENTER,bold=(val==row[0]))
# Chai paper ECs for comparison
T(s,"Paper (Chai et al.) EC values for comparison:",
  Inches(0.4),Inches(5.2),Inches(9),Inches(0.38),sz=13,bold=True,c=ORANGE)
paper_ecs=[("EC1","6.06","12.50","3","Op"),("EC2","6.80","12.70","6.89","Op"),
           ("EC3","7.28","12.34","10.59","Op"),("EC4","10.46","14.12","25","Op"),
           ("EC5","14.33","14.78","57.64","Parked")]
for i,ec in enumerate(paper_ecs):
    T(s,f"{ec[0]}: Hs={ec[1]}m · Tp={ec[2]}s · Vhub={ec[3]}m/s · {ec[4]}",
      Inches(0.5+i*2.55),Inches(5.65),Inches(2.5),Inches(0.5),sz=11,c=LITE)
BL(s,[
  "EC4 (our) / EC5 (paper): parked turbine in 100-year wind — drives the LARGEST tower loads",
  "10 stochastic realizations per EC  ·  TMax = 4200 s  (600 s transient + 3600 s analysis)",
],Inches(0.4),Inches(6.35),Inches(12.5),Inches(0.95),sz=13)

# ══════════════════════════════════════════════════════════════════════════
# S22  OpenFAST model
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"OpenFAST: IEA 15 MW on UMaine VolturnUS-S Semi-Submersible",
    "Fully coupled aero-hydro-servo-elastic simulation — open-source NREL tool")
BL(s,[
  "Turbine specifications:",
  "   Rated power: 15 MW  ·  Cut-in: 3 m/s  ·  Rated: 10.59 m/s  ·  Cut-out: 25 m/s",
  "   Hub height: 150 m  ·  Rotor diameter: 240 m  ·  Max rotor speed: 7.56 rpm",
  "",
  "Platform: UMaine VolturnUS-S  (4-column semi-submersible steel)",
  "   Mooring: 3 catenary chain lines, 850 m length each, MBS = 22,286 kN",
  "   Water depth: 200 m  ·  Fairlead depth: 14 m",
  "",
  "Physics modules used in this simulation:",
  "   AeroDyn  — blade element momentum + generalised dynamic wake aerodynamics",
  "   ServoDyn — ROSCO controller (blade pitch + generator torque)",
  "   ElastoDyn — structural dynamics: rotor, drivetrain, tower, platform",
  "   HydroDyn  — potential flow (WAMIT data) + 2nd-order wave forces (QTF)",
  "   MoorDyn   — dynamic mooring line simulation",
  "   SeaState  — JONSWAP irregular wave generation (WaveMod=2)",
  "",
  "Validated against: DLW, DIEGO, HAWC2, OpenFAST (Chai et al. Tables 4–7)",
  "   All four tools agree within <10% on tower base shear and bending moment ✓",
],Inches(0.4),Inches(1.15),Inches(12.5),Inches(5.9),sz=14)

# ══════════════════════════════════════════════════════════════════════════
# S23  Simulation pipeline + caching
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"Automated Simulation Pipeline",
    "Python script modifies input files, runs OpenFAST, logs outputs — with caching")
BL(s,[
  "For each EC × realization the script does:",
  "   1.  Check if .outb already exists  →  if yes, SKIP (cache hit)",
  "   2.  Copy the base model directory to a fresh run folder",
  "   3.  Copy InflowWind file from IEA-15-240-RWT/ to run dir (fixes relative path)",
  "   4.  Patch HWindSpeed = V_hub in InflowWind (steady wind at hub height, 150 m)",
  "   5.  Patch WaveHs, WaveTp, WaveTMax, WaveSeed in SeaState  (JONSWAP parameters)",
  "   6.  Patch TMax = 4200 s in .fst",
  "   7.  Call openfast subprocess, capture stdout",
  "   8.  Record output file path to run_log for post-processing",
  "",
  "Run directories placed as siblings of IEA-15-240-RWT-UMaineSemi/",
  "   (required so ../IEA-15-240-RWT/ relative file paths resolve correctly)",
  "",
  "Caching means the cell is SAFE TO RE-RUN after kernel crash or partial completion",
  "   To invalidate and rerun: delete EC*_r* folders manually",
  "",
  "Computation scale:",
  "   4 ECs × 10 realizations = 40 runs  ·  ~70 min each  →  ~47 hrs sequential",
  "   Paper: 5 ECs × 50 realizations = 250 runs  →  ~292 hrs sequential",
  "   With 4 CPU cores in parallel: ~12 hrs for our 40-run set",
],Inches(0.4),Inches(1.15),Inches(12.5),Inches(5.9),sz=14)

# ══════════════════════════════════════════════════════════════════════════
# S24  What outputs we extract
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"Simulation Outputs: What We Extract and Why",
    "Three structural response channels — chosen to match Chai et al. 2024 Table 9–11")
BL(s,[
  "From each .outb file (binary format, read with openfast_io FASTOutputFile):",
  "   Discard first 600 s (transient) → analyse 3600 s per realization",
  "",
  "FAIRTEN1  — fairlead tension in mooring line #1  [N → convert to kN]",
  "   Mooring line #1 faces the wave direction alone",
  "   Lines #2 and #3 are symmetrically distributed → less loaded",
  "   This is the MOST LOADED LINE — critical for mooring integrity",
  "   Safety check: max tension < 60% MBS = 0.6 × 22,286 = 13,372 kN",
  "",
  "TwrBsFyt  — tower base fore-aft shear force  [kN]",
  "   Force at the bottom of the tower in the wave/wind direction",
  "   Key indicator of overturning risk",
  "",
  "TwrBsMyt  — tower base fore-aft bending moment  [kN·m]",
  "   Most demanding for structural assessment (largest lever arm)",
  "   Safety check: must be < LULt_moment (Eurocode 3, S355 steel)",
  "   Paper allowables: shear = 1.59×10⁵ kN · moment = 1.08×10⁶ kN·m",
  "",
  "Power spectral density (PSD) of each channel shows which physical",
  "   phenomena dominate: pitch frequency, wave frequency, 1P/3P blade",
],Inches(0.4),Inches(1.15),Inches(12.5),Inches(5.9),sz=14)

# ══════════════════════════════════════════════════════════════════════════
# S25  Gumbel method — full derivation
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"Extreme Value Method 1 — Gumbel Block Maxima",
    "Simple, transparent, used in the paper — but tends to be conservative")
BL(s,[
  "Intuition: if you run 10 simulations and take the maximum of each,",
  "   you have a sample from the DISTRIBUTION OF MAXIMA.",
  "   The Gumbel distribution is theoretically the limiting distribution of maxima.",
  "",
  "Step 1: extract one block maximum per realization  →  10 values per EC per channel",
  "",
  "Step 2: fit Gumbel CDF by MLE:",
  "   F(x) = exp{ −exp[ −(x−μ)/β ] }",
  "   Two parameters: location μ  (shifts the distribution left/right)",
  "                   scale β   (controls how spread out the tail is)",
  "",
  "Step 3: diagnostic — Gumbel probability paper:",
  "   y = −ln(−ln(F(x)))  plotted vs x  should be a STRAIGHT LINE",
  "   Deviation from linearity → Gumbel is not a good fit",
  "",
  "Step 4: read off extreme value at λ = 0.01 (1% exceedance per 1-h sea state):",
  "   x* = μ − β · ln(−ln(1−0.01))  =  μ + 4.60·β",
  "",
  "Known limitation:",
  "   Uses only 1 value per run → ignores hundreds of peaks in each time series",
  "   This makes it conservative (overestimates the true extreme)",
],Inches(0.4),Inches(1.15),Inches(12.5),Inches(5.9),sz=14)

# ══════════════════════════════════════════════════════════════════════════
# S26  ACER method — full derivation
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"Extreme Value Method 2 — ACER  (Naess & Gaidai 2009)",
    "Average Conditional Exceedance Rate — uses every peak, accounts for wave-to-wave correlation")
BL(s,[
  "Motivation: real ocean waves are correlated — a big wave often follows another big wave.",
  "   If you just count exceedances ignoring this, you double-count dependent events.",
  "   ACER conditions on the PREVIOUS peak to remove this bias.",
  "",
  "Definition: ACER₂(η) = P(Xⱼ > η | Xⱼ₋₁ ≤ η)  averaged over all peaks j",
  "   This is the rate at which peaks exceed threshold η,",
  "   given the immediately preceding peak did NOT exceed η  (k=2 Markov order)",
  "",
  "Empirical estimation from R realizations:",
  "   ACER₂(η) = (1/R) × Σᵣ  #{j≥2: Xⱼ>η, Xⱼ₋₁≤η} / Nᵣ",
  "",
  "Tail fitting (parametric extrapolation):",
  "   ê(η) ≈ q · exp{ −a·(η − b)^c }   for η ≥ η₀",
  "   4 parameters (q, a, b, c) fitted by weighted least-squares in log-space",
  "   η₀ = lower bound of the fitting region (upper 20% of peaks)",
  "",
  "Find extreme value η*:  set  ê(η*) = 0.01  →  solve for η*",
  "",
  "Why it works better than Gumbel:",
  "   Uses EVERY PEAK (hundreds per run vs 1 for Gumbel) → much more data",
  "   Explicitly models peak-to-peak correlation  →  less bias",
  "   Paper: ACER gives lower (less conservative) and more accurate estimates",
],Inches(0.4),Inches(1.15),Inches(12.5),Inches(5.9),sz=14)

# ══════════════════════════════════════════════════════════════════════════
# S27  Gumbel vs ACER table
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"Gumbel vs ACER — Comparison and Paper Findings")
cols_h=["","Gumbel","ACER₂"]
col_w=[Inches(3.4),Inches(4.2),Inches(4.2)]
col_x=[Inches(0.5),Inches(4.0),Inches(8.3)]
for col,x,w in zip(cols_h,col_x,col_w):
    BOX(s,x,Inches(1.1),w,Inches(0.55),fill=MID)
    T(s,col,x+Inches(0.05),Inches(1.13),w-Inches(0.1),Inches(0.5),
      sz=15,bold=True,al=PP_ALIGN.CENTER)
rows_cmp=[
  ("Data used","1 block maximum per run","All local peaks (hundreds per run)"),
  ("Sample size (10 runs)","10 values","~500–2000 peaks total"),
  ("Wave correlation","Ignored","Modelled via Markov k=2 conditioning"),
  ("Tail model","Linear on Gumbel paper","Parametric:  q·exp{−a(η−b)^c}"),
  ("Typical extreme estimate","Higher — conservative","Lower — closer to true value"),
  ("Uncertainty (paper Table 13)","Larger CI","Most values < 0.10 — very tight"),
  ("Best use case","Quick check, few simulations","Accurate estimate, more data available"),
]
for ri,(lbl,g,a) in enumerate(rows_cmp):
    yr=Inches(1.75+ri*0.65)
    rb=RGBColor(0x12,0x35,0x6A) if ri%2==0 else DARK
    for val,x,w in zip([lbl,g,a],col_x,col_w):
        BOX(s,x,yr,w,Inches(0.6),fill=rb,border=MID)
        T(s,val,x+Inches(0.05),yr+Inches(0.06),w-Inches(0.1),Inches(0.54),
          sz=12,c=ORANGE if x==col_x[0] else WHITE,al=PP_ALIGN.CENTER)
T(s,"Paper finding (Chai et al. 2024, Section 6):",
  Inches(0.5),Inches(6.4),Inches(4.5),Inches(0.4),sz=14,bold=True,c=ORANGE)
T(s,"ACER estimates are consistently smaller than Gumbel. Under extreme conditions (EC5), "
   "ACER results are closer to the actual situation and are recommended for short-term prediction.",
  Inches(0.5),Inches(6.8),Inches(12.3),Inches(0.5),sz=13,c=LITE)

# ══════════════════════════════════════════════════════════════════════════
# S28  VTK / ParaView
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"OpenFAST VTK Animation — Watching the Physics",
    "WrVTK=2, VTK_type=1: exports solid mesh frames every timestep for ParaView")
BL(s,[
  "OpenFAST VTK output flags in the .fst file:",
  "   WrVTK = 2    → animation (one .vtp mesh file per frame)",
  "   VTK_type = 1 → solid surfaces (blades, tower cylinder, platform)",
  "   VTK_fps = 15 → 15 frames per second",
  "",
  "Two sets of .vtp files created in vtk/ folder:",
  "   *.ED.*.vtp  — ElastoDyn mesh: tower, blades, nacelle  (structural deformation)",
  "   *.MD.*.vtp  — MoorDyn mesh:   3 mooring lines  (position as they go taut/slack)",
  "",
  "Load in ParaView:",
  "   File → Open → vtk/ folder → select *.ED..vtp (auto-groups as time series)",
  "   Click Apply → Press Play ▶ to animate",
  "   Load both .ED and .MD groups together in the same scene",
  "",
  "What you can observe:",
  "   Platform surge (forward-backward motion with waves)",
  "   Platform pitch (tilting due to wave loading + wind thrust)",
  "   Blade flex (tip deflection under aerodynamic load)",
  "   Mooring lines going taut on upwave side, slack on downwave side",
  "",
  "Used a 60-second test run for visual validation of the simulation setup",
  "   Confirmed: turbine spins up, platform responds to waves, mooring holds",
],Inches(0.4),Inches(1.15),Inches(7.5),Inches(5.9),sz=14)
BOX(s,Inches(7.9),Inches(1.2),Inches(5.2),Inches(5.6),
    fill=RGBColor(0x06,0x16,0x2E),border=MID)
T(s,"[ ParaView 3D View ]\n\nBlades: solid surfaces, flexing\nTower: cylinder, deforming\nPlatform: pitching & surging\nMooring: 3 lines, taut/slack\n\nBackground: ocean-dark blue\nCamera: orbit around turbine",
  Inches(8.0),Inches(2.8),Inches(4.8),Inches(2.8),sz=14,c=GREY,al=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
# S29  PhD context — climate extension
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"PhD Context: How This Feeds the Climate-Change Extension",
    "The replication is Stage 1 — the thesis contribution is the temporal evolution")
BL(s,[
  "What the replication gives us (Stages 1–3):",
  "   A validated pipeline: ERA5 → joint distribution → IFORM contour → OpenFAST → Gumbel/ACER",
  "   Verified against Chai et al. 2024 EC values and structural response trends",
  "",
  "The PhD contribution — Stage 4 (climate non-stationarity):",
  "   Repeat the ENTIRE pipeline using moving 30-year windows across 1950–2025",
  "   e.g., 1950–1980, 1951–1981, ..., 1995–2025  (46 overlapping windows)",
  "   Track how the 100-year return level of (Uw, Hs, Tp) changes over time",
  "   Regenerate ECs per window  →  how do design conditions shift with climate?",
  "   Run reduced OpenFAST campaigns at selected windows (early/mid/recent decade)",
  "   Quantify: how much have extreme STRUCTURAL LOADS changed since 1950?",
  "",
  "Questions to answer for the supervisor now:",
  "   1. Primary study region: Irish Sea first, then extend to Northern North Sea?",
  "   2. 10 m only (paper-consistent) OR carry both 10 m and 100 m as formal sensitivity?",
  "   3. Window strategy: 30-year fixed, or decadal blocks?",
  "   4. How many OpenFAST runs are feasible per window given compute budget?",
],Inches(0.4),Inches(1.15),Inches(12.5),Inches(5.5),sz=14)
# timeline boxes
for i,(yr,lbl,c) in enumerate([
    ("1950–1980","Window 1\nEarly climate",RGBColor(0x1A,0x6E,0xAE)),
    ("1965–1995","Window N/2\nMid climate",RGBColor(0x27,0x97,0x6B)),
    ("1995–2025","Window M\nRecent climate",RGBColor(0xE7,0x4C,0x3C)),
]):
    BOX(s,Inches(1.0+i*3.8),Inches(6.3),Inches(3.3),Inches(0.8),fill=c,border=WHITE)
    T(s,f"{yr}\n{lbl}",Inches(1.05+i*3.8),Inches(6.35),Inches(3.2),Inches(0.72),
      sz=12,bold=True,al=PP_ALIGN.CENTER)
T(s,"→",Inches(4.4),Inches(6.55),Inches(0.35),Inches(0.4),sz=18,bold=True,c=WHITE,al=PP_ALIGN.CENTER)
T(s,"→",Inches(8.2),Inches(6.55),Inches(0.35),Inches(0.4),sz=18,bold=True,c=WHITE,al=PP_ALIGN.CENTER)
T(s,"Trend in 100-year extremes  →  trend in design loads  →  climate-driven structural risk",
  Inches(0.4),Inches(7.12),Inches(12.5),Inches(0.38),sz=13,c=ORANGE,al=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
# S30  Next steps
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
HDR(s,"Next Steps and Open Questions")
BL(s,[
  "This week — complete the replication baseline:",
  "   Run all 40 OpenFAST simulations (currently code is ready, caching in place)",
  "   Apply Gumbel and ACER to all 3 channels × 4 ECs",
  "   Reproduce Gumbel probability paper plots and ACER convergence plots",
  "   Compare extreme value estimates against Chai et al. Tables 12–14",
  "",
  "Safety check (Section 7 of paper):",
  "   Mooring: max tension < 13,372 kN  (60% MBS of 22,286 kN, API-RP-2SK)",
  "   Tower shear: max < 1.59×10⁵ kN   (Eurocode 3, S355 steel)",
  "   Tower moment: max < 1.08×10⁶ kN·m",
  "",
  "Extend to all 6 ERA5 buoy locations (M1–M6):",
  "   Compare how Irish Sea vs Northern North Sea conditions shift extreme loads",
  "",
  "Scale up to paper methodology:",
  "   Increase N_REAL from 10 → 50 realizations per EC",
  "   Run 4 ECs in parallel on 4 CPU cores to reduce wall-clock time",
  "",
  "Begin climate extension:",
  "   Download full ERA5 to 2025, build sliding-window contour pipeline",
],Inches(0.4),Inches(1.15),Inches(7.5),Inches(5.9),sz=14)
BL(s,[
  "Open questions for supervisor:",
  "",
  "1)  Primary region first: Irish Sea (our ERA5 locations) or replicate",
  "    Johannessen's Northern North Sea exactly before extending?",
  "",
  "2)  Carry both 10 m (paper-consistent) AND 100 m extrapolation as",
  "    formal sensitivity from the start, or just 10 m baseline?",
  "",
  "3)  Window strategy for climate extension:",
  "    30-year rolling windows (standard climate normal) vs decadal blocks?",
  "",
  "4)  How many OpenFAST runs per window are feasible?",
  "    Even 10 realizations × 4 ECs × 46 windows = 1840 runs",
],Inches(8.0),Inches(1.15),Inches(5.1),Inches(5.9),sz=13)

prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
